from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd
import matplotlib.pyplot as plt

# =========================
# DADOS (EDITE AQUI)
# =========================

resumo_2025 = {
    "Obitos": 86,
    "Sinistros": 1746,
    "Vias Urbanas Obitos": 55,
    "Rodovias Obitos": 20,
}

resumo_2024 = {
    "Obitos": 108,
    "Sinistros": 3089,
    "Vias Urbanas Obitos": 62,
    "Rodovias Obitos": 31,
}

modal_data = pd.DataFrame({
    "Modal": ["Pedestre", "Bicicleta", "Motocicleta", "Automóvel", "Caminhão", "Outro"],
    "2024": [13, 3, 48, 2, 0, 1],
    "2025": [9, 9, 33, 3, 0, 1],
})

mensal_total = [6, 9, 3, 2, 8, 7, 1, 5, 4, 5, 1, 4]

# =========================
# FUNÇÃO CRIAR GRÁFICO
# =========================

def criar_grafico_modal(df):
    plt.figure()
    x = range(len(df["Modal"]))
    plt.bar(x, df["2024"])
    plt.bar(x, df["2025"], bottom=df["2024"])
    plt.xticks(x, df["Modal"], rotation=45)
    plt.title("Comparativo de Óbitos por Modal")
    plt.tight_layout()
    plt.savefig("grafico_modal.png")
    plt.close()


def criar_grafico_mensal(valores):
    plt.figure()
    plt.plot(valores, marker="o")
    plt.title("Óbitos por Mês 2025")
    plt.xticks(range(12),
               ["Jan", "Fev", "Mar", "Abr", "Mai", "Jun",
                "Jul", "Ago", "Set", "Out", "Nov", "Dez"])
    plt.tight_layout()
    plt.savefig("grafico_mensal.png")
    plt.close()


# =========================
# GERAR GRÁFICOS
# =========================

criar_grafico_modal(modal_data)
criar_grafico_mensal(mensal_total)

# =========================
# CRIAR APRESENTAÇÃO
# =========================

prs = Presentation()

def slide_titulo(titulo, subtitulo=""):
    slide_layout = prs.slides.add_slide(prs.slide_layouts[1])
    slide = slide_layout
    slide.shapes.title.text = titulo
    slide.placeholders[1].text = subtitulo
    return slide


# =========================
# CAPA
# =========================

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Balanço de Sinistros e Óbitos no Trânsito 2025"
slide.placeholders[1].text = "Ribeirão Preto\nDepartamento de Segurança Viária"


# =========================
# INTRODUÇÃO
# =========================

slide = slide_titulo(
    "Introdução",
    "O relatório apresenta a evolução dos sinistros e óbitos no trânsito, "
    "com foco na redução de vítimas fatais e apoio à tomada de decisão estratégica."
)


# =========================
# QUADRO RESUMO
# =========================

slide = slide_titulo("Quadro Resumo 2024 x 2025")

content = f"""
2025
Óbitos: {resumo_2025['Obitos']}
Sinistros não fatais: {resumo_2025['Sinistros']}

2024
Óbitos: {resumo_2024['Obitos']}
Sinistros não fatais: {resumo_2024['Sinistros']}
"""

slide.placeholders[1].text = content


# =========================
# COMPARATIVO MODAL
# =========================

slide = slide_titulo("Comparativo por Modal")
slide.shapes.add_picture("grafico_modal.png", Inches(1), Inches(2), width=Inches(6))


# =========================
# DETALHAMENTO MENSAL
# =========================

slide = slide_titulo("Óbitos por Mês 2025")
slide.shapes.add_picture("grafico_mensal.png", Inches(1), Inches(2), width=Inches(6))


# =========================
# RANKING
# =========================

slide = slide_titulo(
    "Ranking de Vias Críticas",
    "1º Av. Dr. Francisco Junqueira\n"
    "2º Av. Independência\n"
    "3º Av. Presidente Vargas\n"
    "4º Av. Prof. João Fiúsa\n"
    "5º Av. Maurílio Biagi"
)


# =========================
# SALVAR
# =========================

prs.save("Balanco_Transito_2025.pptx")

print("Apresentação criada com sucesso!")